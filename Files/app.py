# ----------------------------------------
# PDF Chord and Lyric to PPTX Converter
# ----------------------------------------

# ----------------------------------------
# ----------------- Setup ----------------
# -----------------------------------------
# Install dependencies in Terminal window using this command:
# pip3 install streamlit pdfplumber python-pptx

import streamlit as st
import pdfplumber
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import zipfile
import time
import hashlib

# ----------------------------------------
# --------------- Utilities --------------
# ----------------------------------------
CHORD_PATTERN = r"^[A-G](#|b)?(m|maj|min|sus|dim|aug)?\d*(\/[A-G](#|b)?)?$"

# For chord line detection (is this line mostly chords?)
# Accomodate for chord formats like "F#m", "Bbmaj7", "D/F#", etc
CHORD_LINE_REGEX = re.compile(rf"^({CHORD_PATTERN}\s*)+$")

# For transposing chords
CHORD_TOKEN_REGEX = re.compile(
    r"(?<![/A-Za-z])([A-G](?:#|b)?(?:m|maj|min|sus|dim|aug|add)?\d*(?:/[A-G](?:#|b)?)?)"
)

# For transposing chords (only need notes not entire chord)
ROOT_REGEX = r"^([A-G](?:#|b)?)"

NOTES = ['C', 'C#', 'D', 'D#', 'E', 'F',
         'F#', 'G', 'G#', 'A', 'A#', 'B']

# Flat conversion map
FLAT_MAP = {'Db': 'C#', 
            'Eb': 'D#',
            'Gb': 'F#',
            'Ab': 'G#',
            'Bb': 'A#'}

# Enable display preference for flats
DISPLAY_FLATS = {'C#': 'Db',
                 'D#': 'Eb',
                 'F#': 'Gb',
                 'G#': 'Ab',
                 'A#': 'Bb'}

# Available key options for dropdown menu
KEY_OPTIONS = ["C", "C#", "Db", "D", "D#", "Eb", "E", "F",
               "F#", "Gb", "G", "G#", "Ab", "A", "A#", "Bb", "B"]

SECTION_HEADERS = ("INTRO", "VERSE", "CHORUS", "BRIDGE", "TAG", "ENDING",
                   "REFRAIN", "INSTRUMENTAL", "INTERLUDE", "VAMP", "BREAKDOWN",
                   "TURNAROUND", "PRE-CHORUS", "POST-CHORUS", "OUTRO")

# Accommodate for formats like "Verse 1", "Chorus 2", etc.
SECTION_REGEX = re.compile(rf"^({'|'.join(SECTION_HEADERS)})(\s*\d+)?$", re.IGNORECASE)

# Text colors for PPTX slide
COLOR_SECTION = RGBColor(0x6E, 0xC1, 0x38)  # green
COLOR_CHORD   = RGBColor(0xE2, 0xB8, 0x01)  # yellow/gold
COLOR_LYRIC   = RGBColor(0xFF, 0xFF, 0xFF)  # white
COLOR_NOTE    = RGBColor(0xAA, 0xAA, 0xAA)  # gray

# --------------- Functions ---------------
def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    return text

def load_new_text(new_text, source_id):
    # Use case: Refresh textboxes when new upload or new pasted text
    if source_id != st.session_state.get("last_source_id"):
        st.session_state.text_to_use = new_text
        st.session_state.history = [new_text]
        st.session_state.history_index = 0
        st.session_state.last_source_id = source_id

def normalize_pdf_text(text):
    # Normalize line endings
    text = text.replace("\r", "\n")

    # Remove page numbers
    text = re.sub(r"(?m)^\s*\d+\s*$\n?", "", text)

    # Remove dot-only lines
    text = re.sub(r"(?m)^[ \t]*[\.·•]+[ \t]*$", "", text)

    # Replace leading dots with spaces
    text = re.sub(
        r"(?m)^(\s*)([\.]+)",
        lambda m: m.group(1) + " " * len(m.group(2)),
        text
    )

    # Replace other dot runs with spaces
    text = re.sub(
        r"[\.·•]+",
        lambda m: " " * len(m.group()),
        text
    )

    # Normalize unicode spaces
    text = text.replace("\u00A0", " ")

    # Trim trailing spaces
    text = re.sub(r"[ \t]+$", "", text, flags=re.MULTILINE)

    # Clean excessive blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text

def transpose_text(text, steps, use_flats=False):
    lines = text.split("\n")
    result = []
    for line in lines:
        # Only transpose chord lines
        if detect_chord_line(line):
            
            def transpose_match(match):
                # Original chord text
                chord = match.group()

                # Split slash chords
                if "/" in chord:
                    main, bass = chord.split("/", 1)
                else:
                    main = chord
                    bass = None

                # Extract root + suffix
                root_match = re.match(r'^([A-G](?:#|b)?)(.*)$', main)

                if not root_match:
                    return chord
                root = root_match.group(1)
                suffix = root_match.group(2)

                # Normalize flats
                root = FLAT_MAP.get(root, root)

                # Transpose root
                if root in NOTES:
                    idx = NOTES.index(root)
                    new_root = NOTES[(idx + steps) % 12]
                else:
                    new_root = root

                new_chord = new_root + suffix

                # Transpose bass note
                if bass:
                    bass = FLAT_MAP.get(bass, bass)
                    if bass in NOTES:
                        bass_idx = NOTES.index(bass)
                        new_bass = NOTES[(bass_idx + steps) % 12]
                    else:
                        new_bass = bass

                    new_chord += "/" + new_bass

                return new_chord

            line = CHORD_TOKEN_REGEX.sub(transpose_match, line)
        
        result.append(line)

    # Rebuild text 
    result_text = "\n".join(result)

    # Convert sharps -> flats for display
    if use_flats:
        for sharp, flat in DISPLAY_FLATS.items():
            result_text = result_text.replace(sharp, flat)

    return result_text

def detect_chord_line(text):
    # Contains bar notation
    if "|" in text.strip():
        return True

    # Contains many chord-like tokens (D, G2, Asus4, F#, Bm7, etc.)
    tokens = text.strip().split()
    chord_like = sum(bool(re.match(CHORD_PATTERN, t)) for t in tokens)

    return chord_like >= max(1, len(tokens) // 2) 

def detect_sections(text):
    lines = text.split("\n")
    structured = []
    current_section = None

    for line in lines:
        line_clean = re.sub(r"[\.·•]+", " ", line)

        match = re.match(r"^\[?(%s)(\s*\d+)?\b(.*)$" % "|".join(SECTION_HEADERS), line_clean.strip(), re.IGNORECASE)
        if match:
            if current_section:
                structured.append(current_section)
            
            section_type = match.group(1).upper()
            section_num = (match.group(2) or "").strip()
            section_remainder = match.group(3).strip()

            # Normalize header remainder
            section_remainder = re.sub(r"\s*[:\.]+\s*", ":", section_remainder, count=1)
            
            # Extract MD notes by first colon
            if ":" in section_remainder:
                parts = section_remainder.split(":", 1)
                notes_raw = parts[1].strip()
            else:
                notes_raw = ""
            
            # Clean notes
            notes = notes_raw.strip()
            notes = re.sub(r"\s+", " ", notes)

            current_section = {
                "label": f"{section_type} {section_num}".strip(),
                "notes": notes,
                "lines": []
            }
        else:
            if not current_section:
                current_section = {
                    "label": "VERSE",
                    "notes": "",
                    "lines": []
                }
            current_section["lines"].append(line)

    if current_section:
        structured.append(current_section)

    return structured

def load_css():
    st.html("""
    <style>
    /* App background */
    .stApp {
        background-color: #0E0E0E;
        color: #FFFFFF;
    }

    /* Section headers */
    h1, h2, h3 {
        font-family: -apple-system, BlinkMacSystemFont, sans-serif;
        letter-spacing: -0.5px;
    }

    /* Buttons */
    button[kind="primary"] {
        background-color: #6EC138 !important;
        color: black !important;
        border-radius: 8px;
        font-weight: 600;
    }

    /* File uploader */
    section[data-testid="stFileUploader"] {
        border: 1px dashed #444;
        border-radius: 10px;
        padding: 10px;
    }        

    textarea {
        font-family: "Courier New", monospace !important;
        font-size: 16px !important;
        line-height: 1.4 !important;
    }

    .preview-container {
        overflow-y: auto;
        padding-right: 10px;
        border: 1px solid #333;
        border-radius: 8px;
    }

    .preview-container::-webkit-scrollbar {
        width: 6px;
    }
                
    .preview-container::-webkit-scrollbar-thumb {
        background: #555;
        border-radius: 4px;
    }
                
    .slide-separator {
        text-align: center;
        color: #888;
        font-size: 12px;
        margin: 20px 0 10px 0;
    }

    .slide-box {
        background-color: black;
        font-size: 16px;
        font-family: monospace;
        white-space: pre;
        padding: 20px;
        margin-bottom: 10px;
        border-radius: 6px;
        line-height: 1.3;
        margin: 0; 
    }
                
    .section {
        display: block;
        text-align: left !important;
    }
                
    .section { color: #6EC138; }    
    .chord   { color: #E2B801; }    
    .note    { color: #AAAAAA; font-style: italic; }    
    .lyric   { color: #FFFFFF; }    
            
    </style>
    """)

def split_inline(lines):
    result = []
    i = 0

    while i < len(lines):
        line = lines[i]

        # --- Case 1: Chord + lyric pair ---
        if detect_chord_line(line) and i + 1 < len(lines):
            chord_line = line
            lyric_line = lines[i + 1]

            if "||" in chord_line or "||" in lyric_line:
                source = chord_line if "||" in chord_line else lyric_line

                # Find split indices
                split_indices = []
                idx = 0
                while "||" in source[idx:]:
                    pos = source.index("||", idx)
                    split_indices.append(pos)
                    idx = pos + 2

                # Remove markers
                chord_clean = chord_line.replace("||", "")
                lyric_clean = lyric_line.replace("||", "")

                prev = 0
                offset = 0

                for pos in split_indices:
                    adj = pos - offset

                    c_part = chord_clean[prev:adj]
                    l_part = lyric_clean[prev:adj]

                    # Preserve trailing bar if needed
                    c_part = c_part.rstrip()
                    if "|" in c_part and not c_part.endswith("|"):
                        c_part += "|"

                    result.append(c_part)
                    result.append(l_part)

                    prev = adj
                    offset += 2

                # Final segment
                c_part = chord_clean[prev:]
                l_part = lyric_clean[prev:]

                c_part = c_part.rstrip()
                if "|" in c_part and not c_part.endswith("|"):
                    c_part += "|"

                result.append(c_part)
                result.append(l_part)

            else:
                result.append(chord_line)
                result.append(lyric_line)

            i += 2

        # --- Case 2: Single line ---
        else:
            if "||" in line:
                parts = line.split("||")

                for p in parts:
                    p_clean = p.rstrip()

                    if "|" in p_clean and not p_clean.endswith("|"):
                        p_clean += "|"

                    result.append(p_clean)
            else:
                result.append(line)

            i += 1

    return result

def split_slides(text):
    # Normalize trailing spaces
    text = re.sub(r"[ \t]+$", "", text, flags=re.MULTILINE)

    # Split on --- lines
    slides = re.split(r"(?m)^\s*---\s*$", text)

    return [s.strip() for s in slides if s.strip()]

def format_slides(text):
    lines = text.split("\n")
    cleaned_lines = []

    for line in lines:
        # Split indentation from content
        indent = re.match(r"^(\s*)", line).group(1)
        content = line[len(indent):]

        # Section headers
        if re.match(rf"^\s*({'|'.join(SECTION_HEADERS)})(\s*\d+)?\s*:?", content.strip(), re.IGNORECASE):
            formatted_line = indent + f"<span class='section'>{content}</span>"
        
        # Chord lines
        elif detect_chord_line(content):
            formatted_line = indent + f"<span class='chord'>{content}</span>"

        # Lyrics 
        else:
            formatted_line = indent + f"<span class='lyric'>{content}</span>"

        # Notes (apply last to account for inline notes)
        formatted_line = re.sub(
            r"\(.*?\)",
            lambda m: f"<span class='note'>{m.group()}</span>",
            formatted_line
        )

        # Append lines to list
        cleaned_lines.append(formatted_line)

    return "\n".join(cleaned_lines)

def create_ppt(slides):
    prs = Presentation()

    for slide_text in slides:
        # Use a blank slide (6)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Black background to match preview
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)

        # Manually add a textbox
        textbox = slide.shapes.add_textbox(
            Inches(0.5),  # left
            Inches(1),    # top
            Inches(9),    # width
            Inches(5)     # height
        )

        # Format paragraphs
        text_frame = textbox.text_frame
        text_frame.clear()

        for i, line in enumerate(slide_text.split("\n")):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.2
            p._element.get_or_add_pPr().remove_all('a:buChar')

            # Determine base color for this line
            content = line.strip()
            if re.match(rf"^\s*({'|'.join(SECTION_HEADERS)})(\s*\d+)?\s*:?", content, re.IGNORECASE):
                base_color = COLOR_SECTION
            elif detect_chord_line(content) and content:
                base_color = COLOR_CHORD
            else:
                base_color = COLOR_LYRIC

            # Split line into regular segments and parenthetical notes
            segments = re.split(r'(\(.*?\))', line)
            for segment in segments:
                if not segment:
                    continue
                run = p.add_run()
                run.text = segment
                run.font.size = Pt(30)
                run.font.name = "Menlo"
                if segment.startswith('(') and segment.endswith(')'):
                    run.font.color.rgb = COLOR_NOTE
                    run.font.italic = True
                else:
                    run.font.color.rgb = base_color

        # Enable auto-fit
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Add margins to prevent edge overflow
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)

    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    return ppt_buffer

def create_txt(text):
    # Create text file of Edited Text
    return text.encode("utf-8")

# ----------------------------------------
# ---------- Back-end Processes ----------
# ----------------------------------------
if "history" not in st.session_state:
    st.session_state.history = []

if "history_index" not in st.session_state:
    st.session_state.history_index = -1

if "last_edit_time" not in st.session_state:
    st.session_state.last_edit_time = 0

if "raw_text" not in st.session_state:
    st.session_state.raw_text = ""

if "text_to_use" not in st.session_state:
    st.session_state.text_to_use = ""

if "processed_text" not in st.session_state:
    st.session_state.processed_text = ""

if "target_key" not in st.session_state:
    st.session_state.target_key = NOTES[0]  # safe default

if "prev_original_key" not in st.session_state:
    st.session_state.prev_original_key = None


# ----------------------------------------
# --------------- UI / App ---------------
# ----------------------------------------
# Configure the page
st.set_page_config(
    page_title="Chord Lyrics Slide Converter",
    page_icon="🎵",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Enforce dark mode with custom CSS
st.markdown("""
<style>

/* Force dark mode at browser level */
:root {
    color-scheme: dark;
}

/* Base styling */
html, body, [data-testid="stAppViewContainer"] {
    background-color: #0e1117;
    color: #ffffff;
}
            
/* Labels + radio buttons */
label, div[role="radiogroup"] * {
    color: #ffffff !important;
}
            
/* Inputs + buttons + dropdown+ selectboxes */
input, textarea, button, div[data-baseweb="popover"] , ul[role="listbox"],  li[role="option"], div[data-baseweb="select"] > div{
    background-color: #262730 !important;
    color: #ffffff !important;
}

/* Dropdown background */
div[data-baseweb="popover"] {
    background-color: #262730 !important;
}
            
/* Dropdown hover state */
li[role="option"]:hover {
    background-color: #3a3b44 !important;
}

/* Selectbox text + arrow */
div[data-baseweb="select"] span,
div[data-baseweb="select"] svg {
    color: #ffffff !important;
    fill: #ffffff !important;
}
                        
/* Placeholder text */
::placeholder {
    color: #aaaaaa !important;
}

</style>
""", unsafe_allow_html=True)

# Load CSS
load_css()
    
# Load main UI
st.markdown("""
<h1 style='margin-bottom: 0;'>🎵 Chord Lyrics Slide Converter</h1>
<p style='margin-top: 4px; color: #AAAAAA; font-size: 16px;'>
Convert chord charts into clean, presentation-ready slides.
</p>
""", unsafe_allow_html=True)

# Hide theme toggle
st.markdown("""
<style>
/* Hide top-right menu (contains theme toggle) */
#MainMenu {visibility: hidden;}
header {visibility: hidden;}
</style>
""", unsafe_allow_html=True)

# ---------------
# Step 1: Upload
# ---------------
st.header("Step 1 — Upload")
raw_text = ""

# Create toggle
input_mode = st.radio("Choose Input Method", ["Upload PDF", "Paste Text"], horizontal=True)

# Reset text and history when switching input modes
if "last_mode" not in st.session_state:
    st.session_state.last_mode = input_mode

if st.session_state.last_mode != input_mode:
    st.session_state.last_mode = input_mode
    st.session_state.text_to_use = ""
    st.session_state.pop("uploaded_file", None)
    st.session_state.history = []
    st.session_state.history_index = -1
    st.session_state.last_source_id = None
    st.session_state.last_mode = input_mode

if input_mode == "Upload PDF":
    # Option 1: Upload file
    uploaded_file = st.file_uploader("Upload PDF Here", type=["pdf"], key="uploaded_file")
    if uploaded_file:
        raw_text = extract_text_from_pdf(uploaded_file)
        source_id = uploaded_file.name + str(uploaded_file.size)
        load_new_text(normalize_pdf_text(raw_text), source_id)

elif input_mode == "Paste Text":
    # Option 2: Paste text
    pasted_text = st.text_area("Paste Text Here")

    if pasted_text and pasted_text.strip():
       # Avoid reprocessing on every keystroke
       if st.session_state.get("last_text") != pasted_text:
            st.session_state["last_text"] = pasted_text

            # Use hash for stronger detection (when user uploads file with same name)
            paste_id = hashlib.md5(pasted_text.encode()).hexdigest()
            raw_text = pasted_text
            load_new_text(normalize_pdf_text(pasted_text), paste_id)    


# ---------------------
# Step 2: Review & Fix
# ---------------------
if st.session_state.text_to_use:
    st.header("Step 2 — Review & Fix")

    # --- Key transpose section ---
    col1, col2, col3 = st.columns([1, 0.3, 1])

    with col1:
        original_key = st.selectbox("Original Key", ["Select key..."] + KEY_OPTIONS, key="original_key")

        # Auto-sync target_key if transpose key not selected
        if original_key != "Select key...":
            if (
                st.session_state.target_key == st.session_state.prev_original_key
                or st.session_state.prev_original_key is None
            ):
                st.session_state.target_key = original_key
        st.session_state.prev_original_key = original_key

    with col2:
        st.html("<div style='text-align:center;padding-top:30px;'>→</div>")

    with col3:
        target_key = st.selectbox("Transpose To", ["Select key..."] + KEY_OPTIONS, key="target_key", disabled=(original_key == "Select key..."))

    # Normalize keys to sharps for calculating steps
    orig = FLAT_MAP.get(original_key, original_key)
    targ = FLAT_MAP.get(target_key, target_key)
    steps = 0
    if orig in NOTES and targ in NOTES:
        steps = (NOTES.index(targ) - NOTES.index(orig)) % 12

    # --- Toolbar + Textboxes ---
    col_left, col_right = st.columns([1, 1])
    with col_left:
        # --- Toolbar section
        toolbar_col1, toolbar_col2, toolbar_col3 = st.columns([4, 2, 2])
        with toolbar_col1:
            with st.popover("❓ Help Menu"):
                st.markdown("### How to Edit Text")
                st.markdown("• Use || for line breaks")
                st.markdown("• Use --- for slide breaks")
                st.markdown("• Press Cmd + Enter to apply")

        with toolbar_col2:
            if st.button("↩️", help="Undo Change", disabled=st.session_state.history_index <= 0):
                st.session_state.history_index -= 1
                st.session_state.text_to_use = st.session_state.history[st.session_state.history_index]
                st.rerun()
        with toolbar_col3:
            if st.button("↪️", help="Redo Change", disabled=st.session_state.history_index >= len(st.session_state.history) - 1):
                st.session_state.history_index += 1
                st.session_state.text_to_use = st.session_state.history[st.session_state.history_index]
                st.rerun()
        
        # --- Editable Text textbox section
        st.subheader("Editable Text")
        
        # Dynamically adjust box height based on content
        lines = st.session_state.text_to_use.count("\n") + 1
        height = min(1000, max(400, lines * 24))
        offset = 26     # to account for the Preview Slides padding  

        # Render Edited Text box
        st.text_area(
            "",
            height=height + offset,
            key="text_to_use",
        )
        text_to_use = st.session_state.text_to_use

        # Store current text and history for undo/redo
        current = st.session_state.text_to_use
        history = st.session_state.history
        index = st.session_state.history_index

        # Initalize first state
        if index == -1:
            st.session_state.history = [current]
            st.session_state.history_index = 0
        
        # Save history only when text changes
        elif current != history[index]:
            now = time.time()

            # debounce (prevents saving every keystroke)
            if now - st.session_state.last_edit_time > 0.4:
                # cut off redo history
                st.session_state.history = history[:index + 1]

                # append new state
                st.session_state.history.append(current)
                st.session_state.history_index += 1

                st.session_state.last_edit_time = now

    # --- Slide Preview section
    with col_right:
        st.subheader("")
        st.subheader("Slide Preview")
        
        text_to_use = st.session_state.text_to_use
        processed_text = text_to_use

        slides = split_slides(text_to_use)
        slides_html= ""

        processed_slides = []

        for i, slide in enumerate(slides):
            # Split into lines
            lines = slide.split("\n")

            # Split lines on ||
            lines = split_inline(lines)

            # Rejoin slide
            processed_slide = "\n".join(lines)

            # Apply transpose
            use_flats = isinstance(target_key, str) and "b" in target_key
            if original_key != "Select key..." and target_key != "Select key...":
                processed_slide = transpose_text(processed_slide, steps, use_flats=use_flats)

            # Store processed slide for export
            processed_slides.append(processed_slide)

            # Format slides
            formatted_slide = format_slides(processed_slide)
            formatted_slide = formatted_slide.replace("\n", "<br>")

            # Create a slide separator
            slides_html += f"""<div class="slide-separator">
            ──────── Slide {i+1} ────────
            </div>
            <div class="slide-box">
            {formatted_slide}
            </div>"""
        
        st.session_state.processed_text = "\n---\n".join(processed_slides)

        # Render Slide Preview
        st.html("<div style='height: 5px;'></div>")
        st.html(f"""<div class="preview-container" style="height:{height}px;">
        {slides_html}
        </div>"""
        )

    # ----------------
    # Step 3: Export
    # ----------------
    st.header("Step 3 — Export")
    
    # --- File name input
    file_name_input = st.text_input("Enter File Name", key="file_name_input")

    # Track confirmed filename
    if "confirmed_file_name" not in st.session_state:
        st.session_state.confirmed_file_name = ""

    # Detect Enter (value change triggers rerun)
    if file_name_input != st.session_state.confirmed_file_name:
        st.session_state.confirmed_file_name = file_name_input.strip()

    # Add reminder to press Enter if needed
    if not st.session_state.confirmed_file_name:
        st.caption("Press Enter to confirm file name")
    else:
        st.caption(f"✅ File name confirmed: {st.session_state.confirmed_file_name}")
    
    file_name = st.session_state.confirmed_file_name
    file_name = file_name.replace(".pptx", "").replace(".txt", "")

    ppt_file = create_ppt(processed_slides)

     # --- Create zip file for Download All button
    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, "w") as zip_file:
        zip_file.writestr(f"{file_name}.pptx", ppt_file.getvalue())
        zip_file.writestr(f"{file_name}.txt", text_to_use.encode("utf-8"))
    zip_buffer.seek(0)

    if file_name:
        col1, col2, col3 = st.columns(3)

        with col1:
            st.download_button(
                label="📦 Download All Files",
                data=zip_buffer,
                file_name=f"{file_name}.zip",
                mime="application/zip"
            )

        with col2:
            st.download_button(
                label="🖥 Download Slides",
                data=ppt_file,
                file_name=f"{file_name}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        with col3:
            st.download_button(
                label="🗒 Download Text File",
                data=st.session_state.get("processed_text") or st.session_state.get("text_to_use", ""),
                file_name=f"{file_name}.txt",
                mime="text/plain"
            )
        
        
    