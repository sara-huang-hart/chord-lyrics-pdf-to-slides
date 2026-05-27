# ----------------------------------------
# PDF Chord and Lyric to PPTX Converter
# ----------------------------------------

# ----------------------------------------
# ----------------- Setup ----------------
# -----------------------------------------
# Install dependencies in Terminal window using this command:
# pip3 install streamlit pdfplumber python-pptx

import streamlit as st
import pdfplumber
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import time

# ----------------------------------------
# --------------- Utilities --------------
# ----------------------------------------
CHORD_PATTERN = r"^[A-G](#|b)?(m|maj|min|sus|dim|aug)?\d*(\/[A-G](#|b)?)?$"

# For chord line detection (is this line mostly chords?)
# Accomodate for chord formats like "F#m", "Bbmaj7", "D/F#", etc
CHORD_LINE_REGEX = re.compile(rf"^({CHORD_PATTERN}\s*)+$")

# For transposing chords
CHORD_TOKEN_REGEX = re.compile(
    r"(?<![/A-Za-z])([A-G](?:#|b)?(?:m|maj|min|sus|dim|aug|add)?\d*(?:/[A-G](?:#|b)?)?)"
)

# For transposing chords (only need notes not entire chord)
ROOT_REGEX = r"^([A-G](?:#|b)?)"

NOTES = ['C', 'C#', 'D', 'D#', 'E', 'F',
         'F#', 'G', 'G#', 'A', 'A#', 'B']

# Flat conversion map
FLAT_MAP = {'Db': 'C#', 
            'Eb': 'D#',
            'Gb': 'F#',
            'Ab': 'G#',
            'Bb': 'A#'}

# Enable display preference for flats
DISPLAY_FLATS = {'C#': 'Db',
                 'D#': 'Eb',
                 'F#': 'Gb',
                 'G#': 'Ab',
                 'A#': 'Bb'}

# Available key options for dropdown menu
KEY_OPTIONS = [
    "Select key...",
    "C", "C#", "Db", "D", "D#", "Eb", "E", "F",
    "F#", "Gb", "G", "G#", "Ab", "A", "A#", "Bb", "B"]

SECTION_HEADERS = ("INTRO", "VERSE", "CHORUS", "BRIDGE", "TAG", "ENDING",
                   "REFRAIN", "INSTRUMENTAL", "INTERLUDE", "VAMP", "BREAKDOWN",
                   "TURNAROUND", "PRE-CHORUS", "PRECHORUS", "POST-CHORUS", "POSTCHORUS", "OUTRO")

# Accommodate for formats like "Verse 1", "Chorus 2", etc.
SECTION_REGEX = re.compile(rf"^({'|'.join(SECTION_HEADERS)})(\s*\d+)?$", re.IGNORECASE)

# Text colors for PPTX slide
COLOR_SECTION = RGBColor(0x6E, 0xC1, 0x38)  # green
COLOR_CHORD   = RGBColor(0xE2, 0xB8, 0x01)  # yellow/gold
COLOR_LYRIC   = RGBColor(0xFF, 0xFF, 0xFF)  # white
COLOR_NOTE    = RGBColor(0xAA, 0xAA, 0xAA)  # gray

# --------------- Functions ---------------
def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    return text

def normalize_pdf_text(text):
    # Normalize line endings
    text = text.replace("\r", "\n")

    # Remove page numbers
    text = re.sub(r"(?m)^\s*\d+\s*$\n?", "", text)

    # Remove dot-only lines
    text = re.sub(r"(?m)^[ \t]*[\.·•]+[ \t]*$", "", text)

    # Replace leading dots with spaces
    text = re.sub(
        r"(?m)^(\s*)([\.]+)",
        lambda m: m.group(1) + " " * len(m.group(2)),
        text
    )

    # Replace other dot runs with spaces
    text = re.sub(
        r"[\.·•]+",
        lambda m: " " * len(m.group()),
        text
    )

    # Normalize unicode spaces
    text = text.replace("\u00A0", " ")

    # Trim trailing spaces
    text = re.sub(r"[ \t]+$", "", text, flags=re.MULTILINE)

    # Clean excessive blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text

def transpose_text(text, steps, use_flats=False):
    lines = text.split("\n")
    result = []
    for line in lines:
        # Only transpose chord lines
        if detect_chord_line(line):
            
            def transpose_match(match):
                # Original chord text
                chord = match.group()


                # Split slash chords
                if "/" in chord:
                    main, bass = chord.split("/", 1)
                else:
                    main = chord
                    bass = None

                # Extract root + suffix
                root_match = re.match(r'^([A-G](?:#|b)?)(.*)$', main)

                if not root_match:
                    return chord
                root = root_match.group(1)
                suffix = root_match.group(2)

                # Normalize flats
                root = FLAT_MAP.get(root, root)

                # Transpose root
                if root in NOTES:
                    idx = NOTES.index(root)
                    new_root = NOTES[(idx + steps) % 12]
                else:
                    new_root = root

                new_chord = new_root + suffix

                # Transpose bass note
                if bass:
                    bass = FLAT_MAP.get(bass, bass)
                    if bass in NOTES:
                        bass_idx = NOTES.index(bass)
                        new_bass = NOTES[(bass_idx + steps) % 12]
                    else:
                        new_bass = bass

                    new_chord += "/" + new_bass

                return new_chord

            line = CHORD_TOKEN_REGEX.sub(transpose_match, line)
        
        result.append(line)

    # Rebuild text 
    result_text = "\n".join(result)

    # Convert sharps -> flats for display
    if use_flats:
        for sharp, flat in DISPLAY_FLATS.items():
            result_text = result_text.replace(sharp, flat)

    return result_text

def detect_chord_line(text):
    # Contains bar notation
    if "|" in text.strip():
        return True

    # Contains many chord-like tokens (D, G2, Asus4, F#, Bm7, etc.)
    tokens = text.strip().split()
    chord_like = sum(bool(re.match(CHORD_PATTERN, t)) for t in tokens)

    return chord_like >= max(1, len(tokens) // 2) 

def detect_sections(text):
    lines = text.split("\n")
    structured = []
    current_section = None

    for line in lines:
        line_clean = re.sub(r"[\.·•]+", " ", line)

        match = re.match(r"^\[?(%s)(\s*\d+)?\b(.*)$" % "|".join(SECTION_HEADERS), line_clean.strip(), re.IGNORECASE)
        if match:
            if current_section:
                structured.append(current_section)
            
            section_type = match.group(1).upper()
            section_num = (match.group(2) or "").strip()
            section_remainder = match.group(3).strip()

            # Normalize header remainder
            section_remainder = re.sub(r"\s*[:\.]+\s*", ":", section_remainder, count=1)
            
            # Extract MD notes by first colon
            if ":" in section_remainder:
                parts = section_remainder.split(":", 1)
                notes_raw = parts[1].strip()
            else:
                notes_raw = ""
            
            # Clean notes
            notes = notes_raw.strip()
            notes = re.sub(r"\s+", " ", notes)

            current_section = {
                "label": f"{section_type} {section_num}".strip(),
                "notes": notes,
                "lines": []
            }
        else:
            if not current_section:
                current_section = {
                    "label": "VERSE",
                    "notes": "",
                    "lines": []
                }
            current_section["lines"].append(line)

    if current_section:
        structured.append(current_section)

    return structured

def load_css():
    st.html("""
    <style>
    /* App background */
    .stApp {
        background-color: #0E0E0E;
        color: #FFFFFF;
    }

    /* Section headers */
    h1, h2, h3 {
        font-family: -apple-system, BlinkMacSystemFont, sans-serif;
        letter-spacing: -0.5px;
    }

    /* Buttons */
    button[kind="primary"] {
        background-color: #6EC138 !important;
        color: black !important;
        border-radius: 8px;
        font-weight: 600;
    }

    /* File uploader */
    section[data-testid="stFileUploader"] {
        border: 1px dashed #444;
        border-radius: 10px;
        padding: 10px;
    }        

    textarea {
        font-family: "Courier New", monospace !important;
        font-size: 16px !important;
        line-height: 1.4 !important;
    }

    .preview-container {
        overflow-y: auto;
        padding-right: 10px;
        border: 1px solid #333;
        border-radius: 8px;
    }

    .preview-container::-webkit-scrollbar {
        width: 6px;
    }
                
    .preview-container::-webkit-scrollbar-thumb {
        background: #555;
        border-radius: 4px;
    }
                
    .slide-separator {
        text-align: center;
        color: #888;
        font-size: 12px;
        margin: 20px 0 10px 0;
    }

    .slide-box {
        background-color: black;
        font-size: 16px;
        font-family: monospace;
        white-space: pre;
        padding: 20px;
        margin-bottom: 10px;
        border-radius: 6px;
        line-height: 1.3;
        margin: 0; 
    }
                
    .section {
        display: block;
        text-align: left !important;
    }
                
    .section { color: #6EC138; }    
    .chord   { color: #E2B801; }    
    .note    { color: #AAAAAA; font-style: italic; }    
    .lyric   { color: #FFFFFF; }    
            
    </style>
    """)

def split_inline(lines):
    result = []
    i = 0

    while i < len(lines):
        line = lines[i]

        # --- Case 1: Chord + lyric pair ---
        if detect_chord_line(line) and i + 1 < len(lines):
            chord_line = line
            lyric_line = lines[i + 1]

            if "||" in chord_line or "||" in lyric_line:
                source = chord_line if "||" in chord_line else lyric_line

                # Find split indices
                split_indices = []
                idx = 0
                while "||" in source[idx:]:
                    pos = source.index("||", idx)
                    split_indices.append(pos)
                    idx = pos + 2

                # Remove markers
                chord_clean = chord_line.replace("||", "")
                lyric_clean = lyric_line.replace("||", "")

                prev = 0
                offset = 0

                for pos in split_indices:
                    adj = pos - offset

                    c_part = chord_clean[prev:adj]
                    l_part = lyric_clean[prev:adj]

                    # Preserve trailing bar if needed
                    c_part = c_part.rstrip()
                    if "|" in c_part and not c_part.endswith("|"):
                        c_part += "|"

                    result.append(c_part)
                    result.append(l_part)

                    prev = adj
                    offset += 2

                # Final segment
                c_part = chord_clean[prev:]
                l_part = lyric_clean[prev:]

                c_part = c_part.rstrip()
                if "|" in c_part and not c_part.endswith("|"):
                    c_part += "|"

                result.append(c_part)
                result.append(l_part)

            else:
                result.append(chord_line)
                result.append(lyric_line)

            i += 2

        # --- Case 2: Single line ---
        else:
            if "||" in line:
                parts = line.split("||")

                for p in parts:
                    p_clean = p.rstrip()

                    if "|" in p_clean and not p_clean.endswith("|"):
                        p_clean += "|"

                    result.append(p_clean)
            else:
                result.append(line)

            i += 1

    return result

def split_slides(text):
    # Normalize trailing spaces
    text = re.sub(r"[ \t]+$", "", text, flags=re.MULTILINE)

    # Split on --- lines
    slides = re.split(r"(?m)^\s*---\s*$", text)

    return [s.strip() for s in slides if s.strip()]

def format_slides(text):
    lines = text.split("\n")
    cleaned_lines = []

    for line in lines:
        # Split indentation from content
        indent = re.match(r"^(\s*)", line).group(1)
        content = line[len(indent):]

        # Section headers
        if re.match(rf"^\s*({'|'.join(SECTION_HEADERS)})(\s*\d+)?\s*:?", content.strip(), re.IGNORECASE):
            formatted_line = indent + f"<span class='section'>{content}</span>"
        
        # Chord lines
        elif detect_chord_line(content):
            formatted_line = indent + f"<span class='chord'>{content}</span>"

        # Lyrics 
        else:
            formatted_line = indent + f"<span class='lyric'>{content}</span>"

        # Notes (apply last to account for inline notes)
        formatted_line = re.sub(
            r"\(.*?\)",
            lambda m: f"<span class='note'>{m.group()}</span>",
            formatted_line
        )

        # Append lines to list
        cleaned_lines.append(formatted_line)

    return "\n".join(cleaned_lines)

def create_ppt(slides):
    prs = Presentation()

    for slide_text in slides:
        # Use a blank slide (6)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Black background to match preview
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)

        # Manually add a textbox
        textbox = slide.shapes.add_textbox(
            Inches(0.5),  # left
            Inches(1),    # top
            Inches(9),    # width
            Inches(5)     # height
        )

        # Format paragraphs
        text_frame = textbox.text_frame
        text_frame.clear()

        for i, line in enumerate(slide_text.split("\n")):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.2
            p._element.get_or_add_pPr().remove_all('a:buChar')

            # Determine base color for this line
            content = line.strip()
            if re.match(rf"^\s*({'|'.join(SECTION_HEADERS)})(\s*\d+)?\s*:?", content, re.IGNORECASE):
                base_color = COLOR_SECTION
            elif detect_chord_line(content) and content:
                base_color = COLOR_CHORD
            else:
                base_color = COLOR_LYRIC

            # Split line into regular segments and parenthetical notes
            segments = re.split(r'(\(.*?\))', line)
            for segment in segments:
                if not segment:
                    continue
                run = p.add_run()
                run.text = segment
                run.font.size = Pt(30)
                run.font.name = "Menlo"
                if segment.startswith('(') and segment.endswith(')'):
                    run.font.color.rgb = COLOR_NOTE
                    run.font.italic = True
                else:
                    run.font.color.rgb = base_color

        # Enable auto-fit
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Add margins to prevent edge overflow
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)

    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    return ppt_buffer


# ----------------------------------------
# ---------- Back-end Processes ----------
# ----------------------------------------
if "history" not in st.session_state:
    st.session_state.history = []

if "current_text" not in st.session_state:
    st.session_state.current_text = ""

if "last_committed_text" not in st.session_state:
    st.session_state.last_committed_text = ""

if "redo_stack" not in st.session_state:
    st.session_state.redo_stack = []

# Add keyboard listener for keyboard shortcuts
st.components.v1.html("""
<script>
document.addEventListener('keydown', function(e) {
    // Undo (Cmd/Ctrl + Z)
    if ((e.metaKey || e.ctrlKey) && e.key === 'z') {
        e.preventDefault();
        const input = window.parent.document.querySelector('input[data-testid="undo-trigger"]');
        if (input) {
            input.value = Date.now();
            input.dispatchEvent(new Event('input', { bubbles: true }));
        }
    }

    // Redo (Cmd/Ctrl + Shift + Z)
    if ((e.metaKey || e.ctrlKey) && e.shiftKey && e.key === 'z') {
        e.preventDefault();
        const input = window.parent.document.querySelector('input[data-testid="redo-trigger"]');
        if (input) {
            input.value = Date.now();
            input.dispatchEvent(new Event('input', { bubbles: true }));
        }
    }
});
</script>
""", height=0)


# ----------------------------------------
# --------------- UI / App ---------------
# ----------------------------------------
# Configure the page
st.set_page_config(
    page_title="Chord Lyrics Slide Converter",
    page_icon="🎵",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Load CSS
load_css()
    
# Load main UI
st.markdown("""
<h1 style='margin-bottom: 0;'>🎵 Chord Lyrics Slide Converter</h1>
<p style='margin-top: 4px; color: #AAAAAA; font-size: 16px;'>
Convert chord charts into clean, presentation-ready slides.
</p>
""", unsafe_allow_html=True)


# Hidden triggers input for keyboard shortcuts
undo_trigger = st.text_input("undo_trigger", key="undo_trigger", label_visibility="collapsed")
redo_trigger = st.text_input("redo_trigger", key="redo_trigger", label_visibility="collapsed")

st.markdown("""
    <style>
    /* Hide last 2 text inputs (undo + redo) */
    div[data-testid="stTextInput"]:nth-last-of-type(1),
    div[data-testid="stTextInput"]:nth-last-of-type(2) {
        display: none;
    }
    </style>
    """, unsafe_allow_html=True)

if undo_trigger:
    if st.session_state.history:
        # Push current state to redo stack
        st.session_state.redo_stack.append(st.session_state.current_text)

        # Redo previous
        previous = st.session_state.history.pop()
        st.session_state.current_text = previous
        st.session_state.last_committed_text = previous


if redo_trigger:
    if st.session_state.redo_stack:
        # Save current state to history (so undo still works)
        st.session_state.history.append(st.session_state.current_text)

        # Restore redo state
        next_state = st.session_state.redo_stack.pop()
        st.session_state.current_text = next_state
        st.session_state.last_committed_text = next_state


# ---------------
# Step 1: Upload
# ---------------
st.header("Step 1 — Upload")

uploaded_file = st.file_uploader("Upload PDF", type=["pdf"])
pasted_text = st.text_area("Or paste text")

if uploaded_file:
    raw_text = extract_text_from_pdf(uploaded_file)
elif pasted_text:
    raw_text = pasted_text
else:
    raw_text = ""

# ---------------------
# Step 2: Review & Fix
# ---------------------
if raw_text:
    st.header("Step 2 — Review & Fix")

    # --- Key transpose section ---
    col1, col2, col3 = st.columns([1, 0.3, 1])

    with col1:
        original_key = st.selectbox("Original Key", KEY_OPTIONS, index=0)

    with col2:
        st.html("<div style='text-align:center;padding-top:30px;'>→</div>")

    with col3:
        target_key = st.selectbox("Transpose To", KEY_OPTIONS[1:], disabled=(original_key == "Select key..."))
    
    # Normalize keys for calculation
    if original_key == "Select key...":
        steps = 0
    else:
        calc_original = FLAT_MAP.get(original_key, original_key)
        calc_target = FLAT_MAP.get(target_key, target_key)

        steps = (
            NOTES.index(calc_target)
            - NOTES.index(calc_original)
        ) % 12

    # --- Textboxes section ---
    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Editable Text")

        # Load Help text box
        with st.expander("❓ How to Edit Text"):
            st.markdown("• To insert a line break, type || between the lyrics or chords.")
            st.markdown("• To insert a slide break, type --- between two sections.")
            st.markdown("• Press Cmd + Return to apply changes.")  
            st.markdown( "• To undo a change, use Cmd + Z")
            st.markdown( "• To redo a change, use Cmd + Shift + Z")
            st.markdown("• Zoom out if chords and lyrics are misaligned.")       
        
        # Intialize text box once
        if "edited_text" not in st.session_state:
            st.session_state.edited_text = normalize_pdf_text(raw_text)
        
        # Dynamically adjust box height based on content
        lines = st.session_state.edited_text.count("\n") + 1
        height = min(1000, max(400, lines * 24))
        offset = 26     # to account for the Preview Slides padding  

        # Render Edited Text box
        edited_text = st.text_area(
            "",
            value=st.session_state.edited_text,
            height=height + offset,
            key="editor"
        )

        st.session_state.edited_text = edited_text

        # Capture time to prevent excessive re-renders during typing
        if "last_edit_time" not in st.session_state:
            st.session_state.last_edit_time = 0

        st.session_state.last_edit_time = time.time()

    # -----------------------
    # Step 3: Preview Slides
    # -----------------------
    with col2:
        st.subheader("Slide Preview")
        st.subheader("")
        
        # Get time since last edit
        last_edit_time = st.session_state.get("last_edit_time", 0)
        last_good_text = st.session_state.get("last_good_text", edited_text)

        if time.time() - last_edit_time > 0.3:
            # Only commit if text change is meaningful
            if edited_text != st.session_state.last_committed_text:
                
                # Save previous committed version to history
                if st.session_state.last_committed_text:
                    st.session_state.history.append(st.session_state.last_committed_text)

                # Limit history size to prevent memory issues
                max_history = 50
                if len(st.session_state.history) > max_history:
                    st.session_state.history.pop(0)

                # Clear redo when new edit happens
                st.session_state.redo_stack.clear()

                # Update committed state
                st.session_state.current_text = edited_text
                st.session_state.last_committed_text = edited_text

            text_to_use = st.session_state.current_text

        else:
            # While typing, don't commit yet
            text_to_use = edited_text
           
        slides = split_slides(text_to_use)
        slides_html= ""

        use_flats = "b" in target_key

        processed_slides = []

        for i, slide in enumerate(slides):
            # Split into lines
            lines = slide.split("\n")

            # Split lines on ||
            lines = split_inline(lines)

            # Rejoin slide
            processed_slide = "\n".join(lines)

            # Apply transpose
            processed_slide = transpose_text(processed_slide, steps, use_flats=use_flats)

            # Store processed slide for export
            processed_slides.append(processed_slide)

            # Format slides
            formatted_slide = format_slides(processed_slide)
            formatted_slide = formatted_slide.replace("\n", "<br>")

            # Create a slide separator
            slides_html += f"""<div class="slide-separator">
            ──────── Slide {i+1} ────────
            </div>
            <div class="slide-box">
            {formatted_slide}
            </div>"""

        # Render Slide Preview
        st.html("<div style='height: 5px;'></div>")
        st.html(f"""<div class="preview-container" style="height:{height}px;">
       {slides_html}
        </div>"""
        )

    # Step 4: Export
    st.header("Step 3 — Export")

    ppt_file = create_ppt(processed_slides)

    st.download_button(
        label="⬇️ Download PPTX",
        data=ppt_file,
        file_name="slides.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
